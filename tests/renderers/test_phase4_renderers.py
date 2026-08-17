"""Phase 4 renderer tests: rich-format parser validity (spec section 74).

Every generated file must reopen with its consuming library: docx/xlsx/pptx
via python-docx/openpyxl/python-pptx, eml via the stdlib email parser, PDF
via header/EOF structure, vcf/ics as RFC 6350/5545 text with CRLF endings.

Volume discipline (spec section 72): each file rendered here is well under
1 MiB.
"""

from __future__ import annotations

import random
from email import policy
from email.parser import BytesParser

import pytest

from chaff_generator.content import builders
from chaff_generator.content.context import RenderContext
from chaff_generator.content.template_engine import ChaffTemplateEngine
from chaff_generator.core.hashing import hash_file
from chaff_generator.renderers import build_registry

RICH_FORMATS = ("eml", "docx", "pdf", "xlsx", "pptx", "vcf", "ics")


@pytest.fixture()
def registry():
    return build_registry()


@pytest.fixture()
def world_bank_world(default_bank, world):  # type: ignore[no-untyped-def]
    return default_bank, world


def make_context(world_bank_world, seed: int, desired: int, template_id=None):  # type: ignore[no-untyped-def]
    bank, world = world_bank_world
    rng = random.Random(seed)
    engine = ChaffTemplateEngine(world=world, bank=bank, rng=rng)
    return RenderContext(
        rng=rng,
        world=world,
        bank=bank,
        template_engine=engine,
        desired_size=desired,
        run_id="test",
        app_version="0.1.0",
        template_id=template_id,
    )


# --------------------------------------------------------------------- docx


class TestDocx:
    def test_reopens_with_python_docx(self, world_bank_world, registry, tmp_path):
        bank, _ = world_bank_world
        template = bank.templates().require("prose.quarterly_summary")
        ctx = make_context(world_bank_world, seed=31, desired=48_000, template_id=template.id)
        doc = builders.build_prose_document(template, ctx)
        dest = tmp_path / "out.docx"
        result = registry.get("docx").render(doc, dest, ctx)

        from docx import Document

        opened = Document(str(dest))
        paras = [p.text for p in opened.paragraphs if p.text.strip()]
        assert len(paras) >= 10
        assert opened.core_properties.title == doc.title
        assert opened.core_properties.author == doc.author
        assert dest.stat().st_size == result.size
        assert result.sha256 == hash_file(dest)

    def test_fallback_without_template(self, world_bank_world, registry, tmp_path):
        ctx = make_context(world_bank_world, seed=32, desired=32_000)
        dest = tmp_path / "fallback.docx"
        registry.get("docx").render(None, dest, ctx)
        from docx import Document

        opened = Document(str(dest))
        assert any(p.text.strip() for p in opened.paragraphs)


# ---------------------------------------------------------------------- pdf


class TestPdf:
    def test_header_eof_structure(self, world_bank_world, registry, tmp_path):
        bank, _ = world_bank_world
        template = bank.templates().require("prose.meeting_minutes")
        ctx = make_context(world_bank_world, seed=33, desired=40_000, template_id=template.id)
        doc = builders.build_prose_document(template, ctx)
        dest = tmp_path / "out.pdf"
        registry.get("pdf").render(doc, dest, ctx)

        raw = dest.read_bytes()
        assert raw.startswith(b"%PDF-")
        assert b"%%EOF" in raw[-1024:]

    def test_byte_deterministic(self, world_bank_world, registry, tmp_path):
        """reportlab invariant mode: same seed -> byte-identical PDF."""
        bank, _ = world_bank_world
        template = bank.templates().require("prose.proposal")
        first = tmp_path / "a.pdf"
        second = tmp_path / "b.pdf"
        for seed, dest in ((41, first), (41, second)):
            ctx = make_context(world_bank_world, seed=seed, desired=64_000, template_id=template.id)
            doc = builders.build_prose_document(template, ctx)
            registry.get("pdf").render(doc, dest, ctx)
        assert first.read_bytes() == second.read_bytes()


# --------------------------------------------------------------------- xlsx


class TestXlsx:
    def test_reopens_with_openpyxl(self, world_bank_world, registry, tmp_path):
        bank, _ = world_bank_world
        template = bank.templates().require("tabular.asset_inventory")
        ctx = make_context(world_bank_world, seed=34, desired=48_000, template_id=template.id)
        doc = builders.build_tabular_document(template, ctx)
        dest = tmp_path / "out.xlsx"
        registry.get("xlsx").render(doc, dest, ctx)

        from openpyxl import load_workbook

        wb = load_workbook(dest)
        ws = wb.active
        header = [cell.value for cell in ws[1]]
        assert header == list(doc.sheets[0].columns)
        assert ws.max_row >= 20
        first_header = ws.cell(row=1, column=1)
        assert first_header.font.bold

    def test_currency_and_date_cells_typed(self, world_bank_world, registry, tmp_path):
        bank, _ = world_bank_world
        template = bank.templates().require("tabular.budget_summary")
        ctx = make_context(world_bank_world, seed=35, desired=32_000, template_id=template.id)
        doc = builders.build_tabular_document(template, ctx)
        sheet = doc.sheets[0]
        dest = tmp_path / "typed.xlsx"
        registry.get("xlsx").render(doc, dest, ctx)

        from openpyxl import load_workbook

        ws = load_workbook(dest).active
        if sheet.currency_columns:
            column = sorted(sheet.currency_columns)[0]
            cell = ws.cell(row=2, column=column + 1)
            assert isinstance(cell.value, (int, float)), cell.value
        if sheet.date_columns:
            from datetime import date as date_type

            column = sorted(sheet.date_columns)[0]
            cell = ws.cell(row=2, column=column + 1)
            assert isinstance(cell.value, date_type), cell.value

    def test_fallback_without_template(self, world_bank_world, registry, tmp_path):
        ctx = make_context(world_bank_world, seed=36, desired=16_000)
        dest = tmp_path / "fallback.xlsx"
        registry.get("xlsx").render(None, dest, ctx)
        from openpyxl import load_workbook

        ws = load_workbook(dest).active
        assert ws.max_row >= 5


# --------------------------------------------------------------------- pptx


class TestPptx:
    def test_reopens_with_python_pptx(self, world_bank_world, registry, tmp_path):
        bank, _ = world_bank_world
        template = bank.templates().require("presentation.quarterly_review")
        ctx = make_context(world_bank_world, seed=37, desired=200_000, template_id=template.id)
        doc = builders.build_presentation_document(template, ctx)
        dest = tmp_path / "out.pptx"
        registry.get("pptx").render(doc, dest, ctx)

        from pptx import Presentation

        deck = Presentation(str(dest))
        slides = list(deck.slides)
        assert len(slides) >= len(doc.slides)
        titles = [s.shapes.title.text for s in slides if s.shapes.title is not None]
        assert titles
        assert deck.core_properties.title == doc.title

    def test_large_target_adds_slides(self, world_bank_world, registry, tmp_path):
        bank, _ = world_bank_world
        template = bank.templates().require("presentation.project_kickoff")
        small = make_context(world_bank_world, seed=38, desired=20_000, template_id=template.id)
        big = make_context(world_bank_world, seed=38, desired=400_000, template_id=template.id)
        doc_small = builders.build_presentation_document(template, small)
        doc_big = builders.build_presentation_document(template, big)
        dest_small = tmp_path / "small.pptx"
        dest_big = tmp_path / "big.pptx"
        registry.get("pptx").render(doc_small, dest_small, small)
        registry.get("pptx").render(doc_big, dest_big, big)
        assert dest_big.stat().st_size > dest_small.stat().st_size


# ---------------------------------------------------------------------- eml


class TestEml:
    def _render(self, world_bank_world, registry, dest, seed, template_id, desired=96_000):
        bank, _ = world_bank_world
        template = bank.templates().require(template_id)
        ctx = make_context(world_bank_world, seed=seed, desired=desired, template_id=template.id)
        from chaff_generator.renderers.email import build_email_document

        doc = build_email_document(template, ctx)
        result = registry.get("eml").render(doc, dest, ctx)
        return doc, result, dest

    def test_parses_with_required_headers(self, world_bank_world, registry, tmp_path):
        _, result, dest = self._render(
            world_bank_world,
            registry,
            tmp_path / "out.eml",
            seed=39,
            template_id="email.meeting_invitation",
        )
        with dest.open("rb") as handle:
            msg = BytesParser(policy=policy.default).parse(handle)
        for header in ("From", "To", "Subject", "Date", "Message-ID"):
            assert msg[header], f"missing {header}"
        assert result.sha256 == hash_file(dest)

    def test_multipart_alternative_body(self, world_bank_world, registry, tmp_path):
        _, _, dest = self._render(
            world_bank_world,
            registry,
            tmp_path / "out.eml",
            seed=40,
            template_id="email.newsletter",
            desired=48_000,
        )
        with dest.open("rb") as handle:
            msg = BytesParser(policy=policy.default).parse(handle)
        # With attachments the tree is mixed(alternative(plain, html), att);
        # without, alternative sits at the top. Either way both bodies exist.
        part_types = {part.get_content_type() for part in msg.walk()}
        assert "multipart/alternative" in part_types or "text/plain" in part_types
        plain = msg.get_body(preferencelist=("plain",))
        html = msg.get_body(preferencelist=("html",))
        assert plain is not None and plain.get_content().strip()
        assert html is not None and "<p>" in html.get_content()

    def test_attachment_decodes(self, world_bank_world, registry, tmp_path):
        # A big desired size forces the in-memory attachment path.
        _, _, dest = self._render(
            world_bank_world,
            registry,
            tmp_path / "big.eml",
            seed=42,
            template_id="email.support_reply",
            desired=300_000,
        )
        with dest.open("rb") as handle:
            msg = BytesParser(policy=policy.default).parse(handle)
        attachments = list(msg.iter_attachments())
        assert attachments
        payload = attachments[0].get_content()
        assert isinstance(payload, str) and payload.strip()
        assert attachments[0].get_filename()

    def test_lf_only_line_endings(self, world_bank_world, registry, tmp_path):
        _, _, dest = self._render(
            world_bank_world,
            registry,
            tmp_path / "out.eml",
            seed=43,
            template_id="email.order_confirmation",
        )
        raw = dest.read_bytes()
        assert b"\r" not in raw

    def test_threading_headers_coherent(self, world_bank_world, registry, tmp_path):
        """Any message carrying In-Reply-To must also carry References."""
        for seed in (44, 45, 46, 47, 48):
            _, _, dest = self._render(
                world_bank_world,
                registry,
                tmp_path / f"thread-{seed}.eml",
                seed=seed,
                template_id="email.project_update",
                desired=32_000,
            )
            with dest.open("rb") as handle:
                msg = BytesParser(policy=policy.default).parse(handle)
            if msg["In-Reply-To"]:
                assert msg["References"]
                assert msg["In-Reply-To"] in msg["References"].replace(" ", "").split(",")

    def test_byte_deterministic(self, world_bank_world, registry, tmp_path):
        first = tmp_path / "a.eml"
        second = tmp_path / "b.eml"
        for dest in (first, second):
            self._render(
                world_bank_world,
                registry,
                dest,
                seed=49,
                template_id="email.personal_message",
                desired=48_000,
            )
        assert first.read_bytes() == second.read_bytes()


# ---------------------------------------------------------------------- vcf


class TestVcf:
    def test_crlf_and_structure(self, world_bank_world, registry, tmp_path):
        ctx = make_context(world_bank_world, seed=50, desired=16_000)
        dest = tmp_path / "out.vcf"
        registry.get("vcf").render(None, dest, ctx)

        raw = dest.read_bytes()
        assert raw.endswith(b"\r\n")
        assert b"\n" not in raw.replace(b"\r\n", b"")  # no stray LF
        text = raw.decode()
        assert text.count("BEGIN:VCARD") == text.count("END:VCARD") >= 5
        assert "VERSION:4.0" in text

    def test_cards_have_core_properties(self, world_bank_world, registry, tmp_path):
        ctx = make_context(world_bank_world, seed=51, desired=8_000)
        dest = tmp_path / "out.vcf"
        registry.get("vcf").render(None, dest, ctx)

        text = dest.read_text()
        for card in text.split("END:VCARD"):
            if "BEGIN:VCARD" not in card:
                continue
            for prop in ("FN:", "N:", "EMAIL", "TEL"):
                assert prop in card, f"card missing {prop}"

    def test_unique_people_then_synthesized(self, world_bank_world, registry, tmp_path):
        """World contacts appear once each; no card is repeated verbatim."""
        ctx = make_context(world_bank_world, seed=52, desired=12_000)
        dest = tmp_path / "out.vcf"
        registry.get("vcf").render(None, dest, ctx)

        text = dest.read_text()
        fns = [line for line in text.splitlines() if line.startswith("FN:")]
        assert len(fns) == len(set(fns))

    def test_byte_deterministic(self, world_bank_world, registry, tmp_path):
        digests = []
        for _ in range(2):
            ctx = make_context(world_bank_world, seed=53, desired=9_000)
            dest = tmp_path / "out.vcf"
            result = registry.get("vcf").render(None, dest, ctx)
            digests.append(result.sha256)
        assert digests[0] == digests[1]


# ---------------------------------------------------------------------- ics


class TestIcs:
    def _render(self, world_bank_world, registry, tmp_path, seed, desired=16_000):
        ctx = make_context(world_bank_world, seed=seed, desired=desired)
        dest = tmp_path / "out.ics"
        result = registry.get("ics").render(None, dest, ctx)
        return result, dest

    def test_crlf_and_calendar_structure(self, world_bank_world, registry, tmp_path):
        _, dest = self._render(world_bank_world, registry, tmp_path, seed=54)
        raw = dest.read_bytes()
        assert raw.endswith(b"\r\n")
        assert b"\n" not in raw.replace(b"\r\n", b"")
        text = raw.decode()
        assert text.startswith("BEGIN:VCALENDAR\r\n")
        assert text.rstrip().endswith("END:VCALENDAR")
        assert "VERSION:2.0" in text and "PRODID:" in text
        assert text.count("BEGIN:VEVENT") == text.count("END:VEVENT") >= 5

    def test_events_have_required_properties(self, world_bank_world, registry, tmp_path):
        _, dest = self._render(world_bank_world, registry, tmp_path, seed=55)
        text = dest.read_text()
        for block in text.split("BEGIN:VEVENT")[1:]:
            block = block.split("END:VEVENT")[0]
            for prop in ("UID:", "DTSTAMP:", "DTSTART:", "DTEND:", "SUMMARY:"):
                assert prop in block, f"event missing {prop}"
            start = block.split("DTSTART:")[1].split("\r\n")[0]
            end = block.split("DTEND:")[1].split("\r\n")[0]
            assert start < end  # ISO basic format sorts chronologically

    def test_no_physical_line_exceeds_75_octets(self, world_bank_world, registry, tmp_path):
        _, dest = self._render(world_bank_world, registry, tmp_path, seed=56, desired=32_000)
        for line in dest.read_bytes().split(b"\r\n"):
            assert len(line) <= 75, line[:90]

    def test_folding_unfolds_to_valid_lines(self, world_bank_world, registry, tmp_path):
        """Every continuation line starts with a space; unfolding restores
        logical lines whose property names are valid."""
        _, dest = self._render(world_bank_world, registry, tmp_path, seed=57, desired=24_000)
        physical = dest.read_text().split("\r\n")
        logical: list[str] = []
        for line in physical:
            if line.startswith(" ") and logical:
                logical[-1] += line[1:]
            else:
                logical.append(line)
        known = {
            "BEGIN",
            "END",
            "VERSION",
            "PRODID",
            "CALSCALE",
            "UID",
            "DTSTAMP",
            "DTSTART",
            "DTEND",
            "SUMMARY",
            "LOCATION",
        }
        for line in logical:
            if not line:
                continue
            assert line.split(":", 1)[0].split(";", 1)[0] in known, line[:60]

    def test_uses_world_meeting_dates(self, world_bank_world, registry, tmp_path):
        """Events beyond the world's meeting list fall back to timeline draws,
        but the first events must match real world meeting dates."""
        _, world = world_bank_world
        if not world.meetings:
            pytest.skip("world has no meetings")
        _, dest = self._render(world_bank_world, registry, tmp_path, seed=58)
        text = dest.read_text()
        first_date = world.meetings[0].date
        assert first_date.strftime("%Y%m%d") in text

    def test_byte_deterministic(self, world_bank_world, registry, tmp_path):
        digests = []
        for _ in range(2):
            result, _ = self._render(world_bank_world, registry, tmp_path, seed=59)
            digests.append(result.sha256)
        assert digests[0] == digests[1]


# ------------------------------------------------------- engine integration


class TestEngineRichFormats:
    @pytest.mark.parametrize("fmt", RICH_FORMATS)
    def test_single_format_run_completes(self, fmt, tmp_path):
        """Each rich format flows through the full engine pipeline and every
        produced file is recorded in the manifest with a matching hash."""
        import dataclasses

        from tests.conftest import make_config

        from chaff_generator import ChaffEngine
        from chaff_generator.core.models import FileTypeSetting, TargetMode, TargetSpec

        target = tmp_path / "target"
        target.mkdir()
        config = dataclasses.replace(
            make_config(target),
            target=TargetSpec(path=target, mode=TargetMode.EXACT, amount=384 << 10),
            file_types={fmt: FileTypeSetting(enabled=True)},
        )
        result = ChaffEngine(config).generate()
        assert result.status.value == "completed"
        produced = [p for p in result.run_root.rglob(f"*.{fmt}") if p.is_file()]
        assert produced

        import json

        manifest = json.loads((result.run_root / ".chaff-manifest.json").read_text())
        assert len(manifest["files"]) == len(produced)
        by_path = {f["relative_path"]: f for f in manifest["files"]}
        for path in produced:
            rel = path.relative_to(result.run_root).as_posix()
            assert rel in by_path
            assert by_path[rel]["size"] == path.stat().st_size
            assert by_path[rel]["sha256"] == hash_file(path)
