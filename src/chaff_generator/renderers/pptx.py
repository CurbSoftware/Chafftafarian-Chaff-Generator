"""PPTX renderer — python-pptx over PresentationDocument (spec sections 66, 74).

MVP decks are text and tables only (no charts, per the phase-4 scope);
slides come from the presentation builder, topped up with bullet slides
until the deck approaches the drawn size.
"""

from __future__ import annotations

from typing import TYPE_CHECKING

from chaff_generator.core.hashing import hash_file
from chaff_generator.renderers.base import RendererCapabilities, RenderResult

if TYPE_CHECKING:
    from pathlib import Path

    from chaff_generator.content.context import RenderContext
    from chaff_generator.renderers.base import Renderer
    from chaff_generator.renderers.documents import (
        PresentationDocument,
        SemanticDocument,
    )

CAPABILITIES = RendererCapabilities(
    extension="pptx",
    supports_exact_size=False,
    supports_target_size=True,
    supports_streaming=False,
    semantic_document="presentation",
    size_category="document",
)

#: Decks are XML-heavy: aim at roughly this fraction of the drawn size.
_BODY_FRACTION = 0.12

_LAYOUT_TITLE = 0
_LAYOUT_SECTION = 2
_LAYOUT_BULLETS = 1
_LAYOUT_TABLE = 5


class PptxRenderer:
    id = "pptx"
    capabilities = CAPABILITIES

    def render(
        self,
        document: SemanticDocument | None,
        destination: Path,
        context: RenderContext,
    ) -> RenderResult:
        from pptx import Presentation
        from pptx.util import Inches

        from chaff_generator.renderers.documents import (
            BulletSlide,
            SectionSlide,
            TableSlide,
            TitleSlide,
        )

        if document is None:
            document = self._fallback_document(context)
        from chaff_generator.renderers.documents import PresentationDocument

        assert isinstance(document, PresentationDocument)
        self._top_up(document, context)

        deck = Presentation()
        core = deck.core_properties
        core.title = document.title
        core.author = document.author
        core.last_modified_by = document.author
        core.comments = "Generated synthetic presentation"

        for slide in document.slides:
            if isinstance(slide, TitleSlide):
                page = deck.slides.add_slide(deck.slide_layouts[_LAYOUT_TITLE])
                page.shapes.title.text = slide.title
                page.placeholders[1].text = slide.subtitle
            elif isinstance(slide, SectionSlide):
                page = deck.slides.add_slide(deck.slide_layouts[_LAYOUT_SECTION])
                page.shapes.title.text = slide.title
            elif isinstance(slide, BulletSlide):
                page = deck.slides.add_slide(deck.slide_layouts[_LAYOUT_BULLETS])
                page.shapes.title.text = slide.title
                body = page.placeholders[1].text_frame
                body.text = slide.bullets[0] if slide.bullets else ""
                for bullet in slide.bullets[1:]:
                    paragraph = body.add_paragraph()
                    paragraph.text = bullet
            elif isinstance(slide, TableSlide):
                page = deck.slides.add_slide(deck.slide_layouts[_LAYOUT_TABLE])
                page.shapes.title.text = slide.title
                rows = 1 + len(slide.rows)
                shape = page.shapes.add_table(
                    rows,
                    len(slide.columns),
                    Inches(0.5),
                    Inches(1.6),
                    Inches(9),
                    Inches(0.4 * rows),
                )
                table = shape.table
                for column, name in enumerate(slide.columns):
                    table.cell(0, column).text = str(name)
                for row_index, row in enumerate(slide.rows, start=1):
                    for column, value in enumerate(row):
                        table.cell(row_index, column).text = str(value)

        deck.save(str(destination))
        size = destination.stat().st_size
        return RenderResult(
            path=destination,
            size=size,
            renderer_id=self.id,
            template_id=context.template_id,
            sha256=hash_file(destination),
        )

    # ---------------------------------------------------------------- internals

    def _fallback_document(self, context: RenderContext) -> PresentationDocument:
        from chaff_generator.content import builders
        from chaff_generator.renderers.documents import PresentationDocument

        template_ids = [t.id for t in context.bank.templates().for_kind("presentation")]
        if template_ids:
            template = context.bank.templates().require(template_ids[0])
            return builders.build_presentation_document(template, context)
        return PresentationDocument(
            title="Project Overview",
            author=context.world.primary_user.full_name,
        )

    def _top_up(self, document: PresentationDocument, context: RenderContext) -> None:
        """Add bullet slides until the deck approaches the drawn size."""
        from chaff_generator.content import builders
        from chaff_generator.renderers.documents import BulletSlide

        target_chars = int(context.desired_size * _BODY_FRACTION)
        current = sum(self._slide_chars(slide) for slide in document.slides)
        guard = 0
        while current < target_chars and guard < 2_000:
            count = context.rng.randrange(3, 7)
            bullets = [builders.filler_paragraph(context, sentences=1)[:180] for _ in range(count)]
            slide = BulletSlide(
                title=builders.filler_paragraph(context, sentences=1)[:60].strip().title(),
                bullets=bullets,
            )
            document.slides.append(slide)
            current += self._slide_chars(slide)
            guard += 1

    def _slide_chars(self, slide: object) -> int:
        from chaff_generator.renderers.documents import BulletSlide, TableSlide, TitleSlide

        if isinstance(slide, TitleSlide):
            return len(slide.title) + len(slide.subtitle)
        if isinstance(slide, BulletSlide):
            return len(slide.title) + sum(len(b) for b in slide.bullets)
        if isinstance(slide, TableSlide):
            return (
                len(slide.title)
                + sum(len(str(c)) for c in slide.columns)
                + sum(len(str(cell)) for row in slide.rows for cell in row)
            )
        return len(getattr(slide, "title", ""))


def get_renderer(renderer_id: str) -> Renderer:
    if renderer_id != "pptx":
        raise ValueError(f"pptx module cannot serve renderer id {renderer_id!r}")
    return PptxRenderer()
